from __future__ import annotations

import os
from typing import Any
from pydantic import BaseModel,Field
from langchain_core.tools import StructuredTool
from ..agent.tracker import ActionTracker


class WordDocArgs(BaseModel):
    path:str=Field(...,description="Path to save the .docsx file")
    paragraphs: list[str] = Field(..., description="List of paragraphs to write into the document")

class PptSliceArgs(BaseModel):
    title: str = Field(..., description="Title of the slide")
    content: list[str] = Field(..., description="List of bullet points for the slide")

TEMPLATE_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))), "templete")

def get_available_templates() -> list[str]:
    if os.path.exists(TEMPLATE_DIR):
        return [f for f in os.listdir(TEMPLATE_DIR) if f.endswith(".pptx")]
    return []

# Built-in professional color themes for presentations
PPT_THEMES = {
    "dark_blue":   {"bg": (18, 32, 58),    "accent": (0, 162, 255),  "text": (255, 255, 255), "sub": (180, 210, 255)},
    "dark_green":  {"bg": (15, 40, 30),    "accent": (0, 200, 120),  "text": (255, 255, 255), "sub": (170, 230, 200)},
    "dark_purple": {"bg": (28, 15, 50),    "accent": (180, 80, 255), "text": (255, 255, 255), "sub": (210, 170, 255)},
    "corporate":   {"bg": (240, 244, 248), "accent": (0, 82, 165),   "text": (20, 30, 50),    "sub": (70, 100, 140)},
    "warm":        {"bg": (40, 20, 10),    "accent": (240, 120, 40), "text": (255, 255, 255), "sub": (255, 200, 150)},
    "minimal":     {"bg": (255, 255, 255), "accent": (50, 50, 50),   "text": (30, 30, 30),    "sub": (100, 100, 100)},
    "ocean":       {"bg": (10, 30, 50),    "accent": (0, 190, 210),  "text": (255, 255, 255), "sub": (150, 220, 230)},
    "science":     {"bg": (5, 20, 40),     "accent": (0, 230, 180),  "text": (255, 255, 255), "sub": (160, 230, 210)},
}

class PptDocArgs(BaseModel):
    path: str = Field(..., description="Path to save the .pptx file")
    theme: str = Field(
        "dark_blue",
        description="Color theme for the presentation. Pick based on topic! Options: dark_blue (tech/AI), dark_green (nature/environment), dark_purple (creative/entertainment), corporate (business/finance), warm (marketing/sales), minimal (clean/academic), ocean (science/data), science (research/medicine)"
    )
    slides: list[PptSliceArgs] = Field(..., description="List of slides to create")


class HtmlSlideArgs(BaseModel):
    title: str = Field(..., description="Title of the slide")
    content: list[str] = Field(..., description="List of bullet points for the slide")

class HtmlPresentationArgs(BaseModel):
    path:str=Field(...,description="Path to save the .html presentation file")
    title: str = Field(..., description="Main title of the presentation")
    theme: str = Field(
        "black", 
        description="Reveal.js theme. Choose based on content! Options: black, white, league, sky, beige, simple, serif, blood, night, moon, solarized"
    )
    slides:list[HtmlSlideArgs]=Field(...,description="List of slides to create")

        

def create_office_tools(tracker:ActionTracker)->list[StructuredTool]:

    def list_ppt_templates() -> str:
        """Returns the list of available PowerPoint templates."""
        templates = get_available_templates()
        if not templates:
            return "No templates found. Use 'blank' as the theme."
        return "Available templates: " + ", ".join(templates)

    def create_word_document(path:str,paragraphs:list[str])->str:
        try:
            from docx import Document
        except ImportError:
            return "ERROR:python-docs is not installed".capitalize
        
        try:
            doc=Document()
            for para in paragraphs:
                doc.add_paragraph(para)
            os.makedirs(os.path.dirname(os.path.abspath(path)), exist_ok=True)
            doc.save(path)

            tracker.log(
                action_type="file_create",
                path=path,
                details={"after":f"Word document created with {len(paragraphs)} paragaphs."},
                status="staged"
            )
            return f"SuccessFully created word document at {path}"
        except Exception as e:
            return f"ERROR: Failed to create Word document: {e}"

    def create_powerpoint(path: str, theme: str, slides: list[dict]) -> str:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            from pptx.oxml.ns import qn
            from lxml import etree
        except ImportError:
            return "ERROR: python-pptx is not installed."

        try:
            colors = PPT_THEMES.get(theme, PPT_THEMES["dark_blue"])
            bg_color    = RGBColor(*colors["bg"])
            accent_color = RGBColor(*colors["accent"])
            text_color  = RGBColor(*colors["text"])
            sub_color   = RGBColor(*colors["sub"])

            prs = Presentation()
            prs.slide_width  = Inches(13.33)
            prs.slide_height = Inches(7.5)
            W = prs.slide_width
            H = prs.slide_height
            blank_layout = prs.slide_layouts[6]  # Truly blank layout

            def add_colored_bg(slide, color: RGBColor):
                """Fill the slide background with a solid color."""
                bg = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = color

            def add_accent_bar(slide, color: RGBColor):
                """Add a thin accent bar at the bottom of the slide."""
                bar = slide.shapes.add_shape(
                    1,  # MSO_SHAPE_TYPE.RECTANGLE
                    Inches(0), H - Inches(0.08),
                    W, Inches(0.08)
                )
                bar.fill.solid()
                bar.fill.fore_color.rgb = color
                bar.line.fill.background()

            def add_title_box(slide, title: str, text_color: RGBColor, accent_color: RGBColor):
                """Add a styled title text box."""
                # Accent left bar
                left_bar = slide.shapes.add_shape(
                    1, Inches(0.4), Inches(0.35), Inches(0.07), Inches(0.9)
                )
                left_bar.fill.solid()
                left_bar.fill.fore_color.rgb = accent_color
                left_bar.line.fill.background()

                # Title text
                txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), W - Inches(1.0), Inches(1.1))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = title
                run.font.size = Pt(34)
                run.font.bold = True
                run.font.color.rgb = text_color

            def add_content_box(slide, content_list: list[str], text_color: RGBColor, accent_color: RGBColor):
                """Add formatted bullet points."""
                txBox = slide.shapes.add_textbox(
                    Inches(0.55), Inches(1.55),
                    W - Inches(1.1), H - Inches(2.0)
                )
                tf = txBox.text_frame
                tf.word_wrap = True

                for i, point in enumerate(content_list):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.alignment = PP_ALIGN.LEFT
                    p.space_before = Pt(8)
                    p.space_after = Pt(4)

                    # Bullet symbol run (accent colored)
                    bullet_run = p.add_run()
                    bullet_run.text = ">>  "
                    bullet_run.font.size = Pt(11)
                    bullet_run.font.color.rgb = accent_color
                    bullet_run.font.bold = True

                    # Content run (normal text color)
                    content_run = p.add_run()
                    content_run.text = point
                    content_run.font.size = Pt(19)
                    content_run.font.color.rgb = text_color

            # Build a title slide (first slide)
            title_slide = prs.slides.add_slide(blank_layout)
            add_colored_bg(title_slide, bg_color)
            add_accent_bar(title_slide, accent_color)

            # Centered title for cover slide
            txBox = title_slide.shapes.add_textbox(Inches(1.0), Inches(2.5), W - Inches(2.0), Inches(2.0))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            first_slide_data = slides[0] if slides else {}
            cover_title = first_slide_data.get("title", "") if isinstance(first_slide_data, dict) else getattr(first_slide_data, "title", "")
            run.text = cover_title
            run.font.size = Pt(44)
            run.font.bold = True
            run.font.color.rgb = text_color

            # Accent underline shape below title
            underline = title_slide.shapes.add_shape(
                1, Inches(4.0), Inches(4.7), Inches(5.33), Inches(0.05)
            )
            underline.fill.solid()
            underline.fill.fore_color.rgb = accent_color
            underline.line.fill.background()

            # Content slides
            for slide_data in slides:
                slide = prs.slides.add_slide(blank_layout)
                add_colored_bg(slide, bg_color)
                add_accent_bar(slide, accent_color)

                slide_title = slide_data.get("title", "") if isinstance(slide_data, dict) else getattr(slide_data, "title", "")
                content_list = slide_data.get("content", []) if isinstance(slide_data, dict) else getattr(slide_data, "content", [])

                # Divider line below title area
                divider = slide.shapes.add_shape(
                    1, Inches(0.4), Inches(1.35), W - Inches(0.8), Inches(0.03)
                )
                divider.fill.solid()
                divider.fill.fore_color.rgb = accent_color
                divider.line.fill.background()

                add_title_box(slide, slide_title, text_color, accent_color)
                if content_list:
                    add_content_box(slide, content_list, sub_color, accent_color)

            os.makedirs(os.path.dirname(os.path.abspath(path)), exist_ok=True)
            prs.save(path)

            tracker.log(
                action_type="file_create",
                path=path,
                details={"after": f"PowerPoint created with {len(slides)+1} slides using '{theme}' theme."},
                status="staged"
            )
            return f"Successfully created PowerPoint at {path}"
        except Exception as e:
            return f"ERROR: Failed to create PowerPoint: {e}"

    def create_html_presentation(path:str,title:str,theme:str,slides:list[dict])->str:
        try:
            slides_html=""
            for slide_data in slides:
                slide_title=slide_data.get("title", "") if isinstance(slide_data, dict) else getattr(slide_data, "title", "")
                content_list=slide_data.get("content", []) if isinstance(slide_data, dict) else getattr(slide_data, "content", [])
                list_item="".join([f"<li class='fragment fade-up'>{item}</li>" for item in content_list])
                slide_html=f"""
                <section>
                <h2>{slide_title}</h2>
                <ul>
                    {list_item}
                </ul>
                </section>
                """
                slides_html+=slide_html
            
            html_content=f"""<!doctype html>
            <html lang="en">
            <head>
                <meta charset="utf-8">
                <meta name="viewport" content="width=device-width, initial-scale=1.0, maximum-scale=1.0, user-scalable=no">
                <title>{title}</title>
                <!-- Google Fonts -->
                <link rel="preconnect" href="https://fonts.googleapis.com">
                <link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
                <link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;600;800&display=swap" rel="stylesheet">
                
                <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.3.1/reset.css">
                <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.3.1/reveal.css">
                <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.3.1/theme/{theme}.css" id="theme">
                <style>
                .reveal {{ font-family: 'Inter', sans-serif; }}
                .reveal h1, .reveal h2, .reveal h3 {{ 
                    font-family: 'Inter', sans-serif;
                    text-transform: none; 
                    font-weight: 800;
                    letter-spacing: -0.02em;
                    text-shadow: 2px 2px 4px rgba(0,0,0,0.2);
                }}
                .reveal ul {{ line-height: 1.6; }}
                .reveal li {{ margin-bottom: 0.8em; }}
                </style>
            </head>
            <body>
                <div class="reveal">
                <div class="slides">
                    {slides_html}
                </div>
                </div>
                <script src="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.3.1/reveal.js"></script>
                <script>
                Reveal.initialize({{
                    hash: true,
                    transition: 'slide',
                    backgroundTransition: 'slide',
                    center: true,
                }});
                </script>
            </body>
            </html>"""

            os.makedirs(os.path.dirname(os.path.abspath(path)), exist_ok=True)
            with open(path,"w",encoding="utf-8")as f:
                f.write(html_content)
            
            tracker.log(
                action_type="file_create",
                path=path,
                details={"after":f"Html Presentation created at {path}"},
                status="staged"
            )
            return f"Successfully created HTML presentation at {path}"

        except Exception as e:
            return f"ERROR: Failed to create HTML presentation: {e}"
            

    return[
            StructuredTool.from_function(
                list_ppt_templates,
                name="list_ppt_templates",
                description="CALL THIS FIRST before create_powerpoint. Returns a list of available .pptx template filenames to choose from.",
            ),
            StructuredTool.from_function(
                create_word_document,
                name="create_word_document",
                description="Create or overwrite a Microsoft Word (.docx) document with text paragraphs.",
                args_schema=WordDocArgs,
            ),
            StructuredTool.from_function(
                create_powerpoint,
                name="create_powerpoint",
                description="Create a styled Microsoft PowerPoint (.pptx) presentation using a template. You MUST call `list_ppt_templates` first to get the list of templates, then pass the chosen template filename as `theme`.",
                args_schema=PptDocArgs,
            ),
            StructuredTool.from_function(
                create_html_presentation,
                name="create_html_presentation",
                description="Create an interactive HTML web presentation using Reveal.js. Only use this if the user explicitly asks for a web/HTML presentation.",
                args_schema=HtmlPresentationArgs,
            ),
    ]

                   
                 

    
 
             


        